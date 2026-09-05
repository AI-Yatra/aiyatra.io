#!/usr/bin/env python3
"""Generate two AIYatra PowerPoint decks."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT_DIR = "/Users/moinuddin/Documents/AIYatra/aiyatra.io-website/AIYatra.IO/presentations"
os.makedirs(OUT_DIR, exist_ok=True)

# Brand palette
INK = RGBColor(0x1A, 0x20, 0x2C)
NAVY = RGBColor(0x1E, 0x3A, 0x8A)
DEEP = RGBColor(0x1E, 0x40, 0xAF)
YELLOW = RGBColor(0xFA, 0xCC, 0x15)
YELLOW_DARK = RGBColor(0xCA, 0x8A, 0x04)
CREAM = RGBColor(0xFF, 0xFB, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORAL = RGBColor(0xF4, 0x63, 0x5E)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
SLATE = RGBColor(0x47, 0x5A, 0x6B)
LIGHT_GREY = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x64, 0x74, 0x8B)

WEBSITE = "https://aiyatra.io"
LINKEDIN = "https://www.linkedin.com/company/aiyatra/"
EMAIL = "global.aiyatra@gmail.com"
MEETUP = "https://www.meetup.com/aiyatra/"
GITHUB = "https://github.com/AI-Yatra"
FOOTER_TEXT = f"  {WEBSITE}   •   LinkedIn: linkedin.com/company/aiyatra   •   {EMAIL}   •   meetup.com/aiyatra"

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1.5)
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=INK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(text_frame, text, font_size=14, bold=False, color=INK, alignment=PP_ALIGN.LEFT, space_after=Pt(4), font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    return p

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=INK, bullet_color=None, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(5)
        p.space_before = Pt(1)
        p.line_spacing = line_spacing
        p.level = 0
    return txBox

def add_footer(slide, number_text=None):
    # thin yellow top border
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), Inches(13.33), Pt(4))
    bar.fill.solid(); bar.fill.fore_color.rgb = YELLOW
    bar.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(12.53), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = FOOTER_TEXT
    p.font.size = Pt(7.5)
    p.font.color.rgb = MUTED
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

def add_top_kicker(slide, text, left=Inches(0.6), top=Inches(0.3)):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.2), Inches(0.4))
    box.fill.solid(); box.fill.fore_color.rgb = YELLOW
    box.line.fill.background()
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = box.text_frame.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = INK; p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return box

def add_slide_number(slide, text):
    add_textbox(slide, Inches(12.3), Inches(0.25), Inches(0.6), Inches(0.3), text, font_size=8, color=MUTED, alignment=PP_ALIGN.RIGHT)

# ============================================================
# DECK 1 — AIYatra Journey
# ============================================================
def build_deck1():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ---- 1. Title ----
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    # yellow accent bar top
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Pt(10))
    bar.fill.solid(); bar.fill.fore_color.rgb = YELLOW; bar.line.fill.background()
    add_textbox(s, Inches(0.8), Inches(0.7), Inches(7.5), Inches(0.5), "HYDERABAD  •  FREE, ALWAYS  •  SINCE JUNE 2026", 11, True, YELLOW, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.8), Inches(1.15), Inches(7.5), Inches(2.2), "AIYatra", 72, True, WHITE, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.8), Inches(2.55), Inches(7.5), Inches(0.8), "Research. Build. Transform.", 36, True, YELLOW, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.8), Inches(3.5), Inches(6.8), Inches(1.0), "Democratizing AI Learning — from the first Transformer whiteboard to 2,957 members, 11 hands-on sessions, and the new Student Ambassador Program.", 15, False, WHITE, PP_ALIGN.LEFT)
    # right card
    card = add_rounded_rect(s, Inches(8.4), Inches(0.9), Inches(4.1), Inches(4.9), WHITE)
    add_textbox(s, Inches(8.8), Inches(1.2), Inches(3.3), Inches(0.4), "AT A GLANCE", 10, True, DEEP, PP_ALIGN.LEFT)
    add_textbox(s, Inches(8.8), Inches(1.6), Inches(3.3), Inches(0.6), "2,957  Members", 24, True, INK, PP_ALIGN.LEFT)
    add_textbox(s, Inches(8.8), Inches(2.25), Inches(3.3), Inches(0.4), "11 Sessions  •  4.6 ★ (90 ratings)", 13, True, SLATE, PP_ALIGN.LEFT)
    add_textbox(s, Inches(8.8), Inches(2.85), Inches(3.3), Inches(0.5), "Saturdays @ LSEG, Madhapur, Hyderabad", 12, False, SLATE, PP_ALIGN.LEFT)
    stats = [("1,336", "Total RSVPs"), ("319", "Biggest room"), ("100%", "Free forever")]
    for i, (v, l) in enumerate(stats):
        add_textbox(s, Inches(8.8 + i*1.35), Inches(3.6), Inches(1.3), Inches(0.5), v, 20, True, NAVY, PP_ALIGN.LEFT)
        add_textbox(s, Inches(8.8 + i*1.35), Inches(4.05), Inches(1.3), Inches(0.4), l, 9, True, MUTED, PP_ALIGN.LEFT)
    add_textbox(s, Inches(8.8), Inches(4.7), Inches(3.3), Inches(0.9),
                f"{WEBSITE}\nlinkedin.com/company/aiyatra\n{EMAIL}", 10, False, SLATE, PP_ALIGN.LEFT)
    add_footer(s)

    # ---- 2. What is AIYatra ----
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_top_kicker(s, "✦  What is AIYatra")
    add_slide_number(s, "02")
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.8), "An open door into AI, in the heart of Hyderabad.", 34, True, INK, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(1.7), Inches(7.2), Inches(1.0), "For anyone exploring AI apps, reading Artificial Intelligence: A Modern Approach, building recommenders, or mastering ML with Python. Beginner or expert — you get knowledge sharing, networking, and collaborative projects.", 13.5, False, SLATE, PP_ALIGN.LEFT)
    pillars = [
        ("01  RESEARCH", "Read papers, trace the math, ask naive questions out loud — no gatekeeping, no jargon walls.", NAVY, WHITE),
        ("02  BUILD", "Laptops open, code on screen. Every Saturday ships something real — training loops, agents, demos.", TEAL, WHITE),
        ("03  TRANSFORM", "Skills → careers, side projects → products, strangers → collaborators. That is the yatra.", CORAL, WHITE),
    ]
    for i, (t, b, fill, tc) in enumerate(pillars):
        x = Inches(0.6 + i*4.05)
        card = add_rounded_rect(s, x, Inches(3.1), Inches(3.75), Inches(2.9), WHITE, MUTED, 1)
        # colored header
        hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.1), Inches(3.75), Inches(0.65))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = fill; hdr.line.fill.background()
        hdr.text_frame.word_wrap = True; hdr.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = hdr.text_frame.paragraphs[0]; p.text = t; p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = tc; p.font.name="Calibri"; p.alignment = PP_ALIGN.CENTER
        add_textbox(s, x+Inches(0.3), Inches(3.95), Inches(3.15), Inches(1.8), b, 13, False, SLATE, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(6.25), Inches(12), Inches(0.4), "Led by Super Organizer Khaja Moinuddin Mohammed + volunteer crew  •  Saturdays, mornings IST  •  Bring a laptop, Python 3.10+, curiosity", 10, True, MUTED, PP_ALIGN.LEFT)
    add_footer(s)

    # ---- 3. Timeline ----
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_top_kicker(s, "Journey  •  June → September 2026")
    add_slide_number(s, "03")
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.7), "From a whiteboard to a movement — in 11 weeks.", 30, True, INK, PP_ALIGN.LEFT)
    events = [
        ("JUN 20", "Transformer,\nPaper → Code", "56 attended\nWhere it all began", YELLOW),
        ("JUN 30–\nJUL 14", "Hitchhiker's Guide\nReading #1–4", "25→35→16→40\nOnline evenings", LIGHT_GREY),
        ("JUL 25", "Speculative\nDecoding Lab", "64 attended\n3-hr hands-on lab", LIGHT_GREY),
        ("AUG 01", "Goose AI Agent\nEnd-to-end Demo", "142 attended\nMCP + recipes", LIGHT_GREY),
        ("AUG 08", "DeepSeek-V3\nFrom Scratch", "286 attended\nMLA + MoE + RoPE", YELLOW),
        ("AUG 15", "Linear Algebra\n→ Transformers", "207 attended\nMath by hand", LIGHT_GREY),
        ("AUG 22", "PyTorch\nFoundations", "146 attended\nAutograd by hand", LIGHT_GREY),
        ("SEP 05", "Agentic AI:\nTools That Build Code", "319 going\n100% offline agent", NAVY),
    ]
    for i, (d, t, b, fill) in enumerate(events):
        x = Inches(0.35 + i*1.6)
        is_navy = (fill == NAVY)
        tc = WHITE if is_navy else INK
        sc = WHITE if is_navy else MUTED
        card = add_rounded_rect(s, x, Inches(1.95), Inches(1.5), Inches(4.1), fill, MUTED if not is_navy else NAVY, 1)
        add_textbox(s, x+Inches(0.1), Inches(2.05), Inches(1.3), Inches(0.6), d, 10, True, CORAL if not is_navy else YELLOW, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.1), Inches(2.65), Inches(1.3), Inches(1.0), t, 11, True, tc, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.1), Inches(3.7), Inches(1.3), Inches(1.0), b, 9.5, False, sc if not is_navy else WHITE, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(6.25), Inches(12), Inches(0.4), "★ Aug 5 milestone: 1,000+ members celebrated on LinkedIn  →  2,957 members today  •  Newest chapter: Student Ambassador Program (see slides 11–12)", 10, True, DEEP, PP_ALIGN.LEFT)
    add_footer(s)

    # ---- 4. By the numbers ----
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    add_slide_number(s, "04")
    add_textbox(s, Inches(0.6), Inches(0.4), Inches(5), Inches(0.4), "PROOF, NOT PROMISES", 10, True, YELLOW, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(0.8), Inches(7), Inches(0.8), "The numbers tell the story.", 34, True, WHITE, PP_ALIGN.LEFT)
    nums = [("2,957", "Community members", "Hyderabad + online"), ("11", "Sessions hosted", "Jun 20 → Sep 5"), ("4.6 ★", "Meetup rating", "From 90 ratings"), ("1,336", "Total learners in rooms", "Largest: 319 (Sep 5)")]
    for i, (v, l1, l2) in enumerate(nums):
        x = Inches(0.6 + i*3.05)
        card = add_rounded_rect(s, x, Inches(1.9), Inches(2.85), Inches(2.2), WHITE)
        add_textbox(s, x+Inches(0.25), Inches(2.1), Inches(2.35), Inches(0.7), v, 34, True, NAVY, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.25), Inches(2.85), Inches(2.35), Inches(0.4), l1, 12, True, INK, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.25), Inches(3.25), Inches(2.35), Inches(0.4), l2, 10, False, MUTED, PP_ALIGN.LEFT)
    # top rooms bar
    card = add_rounded_rect(s, Inches(0.6), Inches(4.5), Inches(12.13), Inches(1.7), WHITE)
    add_textbox(s, Inches(1.0), Inches(4.65), Inches(11), Inches(0.4), "BIGGEST ROOMS SO FAR  (Meetup RSVPs at time of publishing)", 9, True, MUTED, PP_ALIGN.LEFT)
    rooms = [("Sep 5 · Agentic AI Build", 319), ("Aug 8 · DeepSeek-V3", 286), ("Aug 15 · Linear Algebra", 207), ("Aug 22 · PyTorch", 146), ("Aug 1 · Goose Demo", 142)]
    maxv = 319
    for i, (label, v) in enumerate(rooms):
        y = Inches(5.05 + i*0.0)  # single row
        x = Inches(1.0 + i*2.35)
        add_textbox(s, x, Inches(5.05), Inches(2.0), Inches(0.3), label, 8.5, True, INK, PP_ALIGN.LEFT)
        # bar bg
        bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.35), Inches(2.0), Pt(14))
        bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT_GREY; bg.line.fill.background()
        fg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.35), Inches(2.0*v/maxv), Pt(14))
        fg.fill.solid(); fg.fill.fore_color.rgb = YELLOW if i==0 else NAVY; fg.line.fill.background()
        add_textbox(s, x+Inches(2.02), Inches(5.32), Inches(0.4), Inches(0.25), str(v), 8.5, True, WHITE, PP_ALIGN.LEFT)
    add_footer(s)

    # ---- 5. Learning tracks ----
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_top_kicker(s, "✦  Curriculum — three tracks")
    add_slide_number(s, "05")
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.7), "One arc: foundations → agents → reading room.", 30, True, INK, PP_ALIGN.LEFT)
    tracks = [
        ("FOUNDATIONS", "Transformer paper → code (Jun 20) • Linear algebra → attention (Aug 15) • PyTorch tensors → training loop (Aug 22)", "Build attention from scratch, Q/K/V by hand, nn.Linear rebuilt. No magic, just math + code.", YELLOW),
        ("AGENTS", "Goose demo (Aug 1) • DeepSeek-V3 live-code (Aug 8) • Agentic coding harness (Sep 5)", "Providers, MCP extensions, MLA/MoE/RoPE, then your own offline coding agent. Zero API keys.", TEAL),
        ("READING ROOM", "Hitchhiker's Guide to Agentic AI — 4 evening sessions, Jun 30 → Jul 14", "Tools, planning, memory, workflows — read together, debated live, ending in a build queue.", CORAL),
    ]
    for i, (t, ev, body, accent) in enumerate(tracks):
        x = Inches(0.6 + i*4.05)
        card = add_rounded_rect(s, x, Inches(1.9), Inches(3.75), Inches(4.2), WHITE, MUTED, 1)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(0.3), Inches(2.15), Pt(18), Pt(18))
        dot.fill.solid(); dot.fill.fore_color.rgb = accent; dot.line.fill.background()
        add_textbox(s, x+Inches(0.65), Inches(2.1), Inches(2.8), Inches(0.35), t, 11, True, INK, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.3), Inches(2.6), Inches(3.15), Inches(1.1), ev, 10.5, True, DEEP, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.3), Inches(3.75), Inches(3.15), Inches(1.8), body, 11.5, False, SLATE, PP_ALIGN.LEFT)
    add_footer(s)

    # ---- 6. Signature sessions table ----
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_top_kicker(s, "Deep dives  •  what we built together")
    add_slide_number(s, "06")
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6), "Every Saturday ships something running on your laptop.", 28, True, INK, PP_ALIGN.LEFT)
    rows = [
        ("Sep 5 · Harnessing Agentic AI", "Agent loop, tools, guardrails — own coding agent, offline", "319"),
        ("Aug 22 · PyTorch Foundations", "Tensors, autograd engine by hand, 5-line training loop", "146"),
        ("Aug 15 · Linear Algebra → Transformers", "Embeddings, Q/K/V, LoRA & SVD — attention in NumPy", "207"),
        ("Aug 8 · DeepSeek-V3 from Scratch", "MLA compression, MoE routing, RoPE scaling in PyTorch", "286"),
        ("Aug 1 · Goose AI Agent Demo", "Providers, context eng., MCP extensions, recipes", "142"),
        ("Jul 25 · Speculative Decoding Lab", "~60 lines PyTorch + HF + llama.cpp, 2–3× faster, lossless", "64"),
        ("Jun 20 · Transformer Paper → Code", "Attention Is All You Need → running encoder block", "56"),
    ]
    # header
    hdr = add_rounded_rect(s, Inches(0.6), Inches(1.75), Inches(12.13), Inches(0.5), NAVY)
    add_textbox(s, Inches(0.9), Inches(1.82), Inches(4.2), Inches(0.35), "SESSION", 9, True, YELLOW, PP_ALIGN.LEFT)
    add_textbox(s, Inches(5.2), Inches(1.82), Inches(5.6), Inches(0.35), "WHAT LAPTOPS LEFT WITH", 9, True, YELLOW, PP_ALIGN.LEFT)
    add_textbox(s, Inches(11.2), Inches(1.82), Inches(1.2), Inches(0.35), "IN ROOM", 9, True, YELLOW, PP_ALIGN.LEFT)
    y = 2.35
    for i, (a, b, c) in enumerate(rows):
        fill = CREAM if i % 2 == 0 else WHITE
        r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(12.13), Inches(0.58))
        r.fill.solid(); r.fill.fore_color.rgb = fill; r.line.fill.background()
        add_textbox(s, Inches(0.9), Inches(y+0.08), Inches(4.2), Inches(0.42), a, 10, True, INK, PP_ALIGN.LEFT)
        add_textbox(s, Inches(5.2), Inches(y+0.08), Inches(5.6), Inches(0.42), b, 10, False, SLATE, PP_ALIGN.LEFT)
        add_textbox(s, Inches(11.2), Inches(y+0.08), Inches(1.2), Inches(0.42), c, 10, True, NAVY, PP_ALIGN.LEFT)
        y += 0.63
    add_footer(s)

    # ---- 7. People + Place ----
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_top_kicker(s, "♥  People who run the yatra")
    add_slide_number(s, "07")
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.7), "A small volunteer crew. A big open room.", 30, True, INK, PP_ALIGN.LEFT)
    people = [
        ("Khaja Moinuddin Mohammed", "Super Organizer & Founder", "Sets the learning arc, hosts Saturdays, keeps the bar high — from linear algebra to agentic harnesses.", YELLOW),
        ("Azeez Syed", "Co-organizer & Host", "Keeps the room running — demos, hands-on labs, Q&A. No learner leaves stuck.", TEAL),
        ("Jagadeeswara Reddy", "Host & Educator", "Turns dense topics into clear, hands-on learning the community can actually use.", CORAL),
    ]
    for i, (n, r, b, accent) in enumerate(people):
        x = Inches(0.6 + i*4.05)
        card = add_rounded_rect(s, x, Inches(1.9), Inches(3.75), Inches(2.9), WHITE, MUTED, 1)
        stripe = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.9), Inches(3.75), Pt(10))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = accent; stripe.line.fill.background()
        add_textbox(s, x+Inches(0.3), Inches(2.25), Inches(3.15), Inches(0.5), n, 13, True, INK, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.3), Inches(2.75), Inches(3.15), Inches(0.35), r.upper(), 9, True, DEEP, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.3), Inches(3.15), Inches(3.15), Inches(1.4), b, 11, False, SLATE, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(5.05), Inches(12.13), Inches(1.0),
                "VENUE  •  LSEG, International Tech Park, Madhapur, Hyderabad  —  Saturdays, mornings IST  •  RSVP on Meetup + Google form at the gate  •  A Saturday: 9:00 doors & check-in → build, not slides → demos, Q&A & hallway track",
                11, False, SLATE, PP_ALIGN.LEFT)
    add_footer(s)

    # ---- 8. Voices ----
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    add_slide_number(s, "08")
    add_textbox(s, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4), "KIND WORDS  •  4.6 ★ FROM 90 RATINGS", 10, True, YELLOW, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7), "Rated by the people who show up.", 32, True, WHITE, PP_ALIGN.LEFT)
    quotes = [
        ("“Walked in knowing nothing about transformers, walked out having built attention from scratch.”", "Priya S. · Data Analyst → ML Engineer"),
        ("“Best Saturday in years. We built a coding agent with no API keys — and it worked.”", "Rahul K. · Backend Engineer"),
        ("“Free, welcoming, genuinely deep — the linear algebra session made math click.”", "Sai Rishita M. · CS Undergraduate"),
        ("“You come for sessions, stay for people. I found my co-founder here.”", "Kiran K. · Founder, AI Startup"),
    ]
    for i, (q, a) in enumerate(quotes):
        x = Inches(0.6 + i*3.05)
        card = add_rounded_rect(s, x, Inches(1.8), Inches(2.85), Inches(3.3), WHITE)
        add_textbox(s, x+Inches(0.25), Inches(2.0), Inches(2.35), Inches(1.8), q, 11, False, INK, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.25), Inches(3.9), Inches(2.35), Inches(0.7), a, 9, True, MUTED, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(5.4), Inches(12.13), Inches(0.5), "Culture: hands-on rhythm — every meetup ends with something running on your laptop  •  Volunteers explain until it clicks", 11, False, WHITE, PP_ALIGN.CENTER)
    add_footer(s)

    # ---- 9. Digital home ----
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_top_kicker(s, "Digital home  •  always open")
    add_slide_number(s, "09")
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.7), "The yatra lives online, too.", 30, True, INK, PP_ALIGN.LEFT)
    cards = [
        ("WEBSITE", "aiyatra.io", "Gatherings, field notes, method, voices, ambassadors + new Blog with 10 meetup seed posts.", YELLOW),
        ("MEETUP", "meetup.com/aiyatra", "Full archive: photos, discussions, ratings, RSVPs — straight from the source.", NAVY),
        ("BLOG + GITHUB", "Blog · github.com/AI-Yatra", "Paper→code write-ups, labs, reading notes. Code & curriculum in the open.", TEAL),
    ]
    for i, (k, url, body, accent) in enumerate(cards):
        x = Inches(0.6 + i*4.05)
        card = add_rounded_rect(s, x, Inches(1.9), Inches(3.75), Inches(3.2), CREAM, MUTED, 1)
        stripe = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Inches(0.3), Inches(2.15), Inches(0.6), Pt(10))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = accent; stripe.line.fill.background()
        add_textbox(s, x+Inches(0.3), Inches(2.45), Inches(3.15), Inches(0.35), k, 10, True, MUTED, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.3), Inches(2.8), Inches(3.15), Inches(0.5), url, 15, True, NAVY if accent!=NAVY else INK, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.3), Inches(3.4), Inches(3.15), Inches(1.3), body, 11, False, SLATE, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(5.4), Inches(12.13), Inches(0.6), "Custom domain + Vercel-ready site  •  Sections: Gatherings · Field Notes · The Yatra Way · Kind Words · Movement · Ambassadors · Blog", 10.5, False, MUTED, PP_ALIGN.LEFT)
    add_footer(s)

    # ---- 10. Ambassador spotlight ----
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Pt(10))
    bar.fill.solid(); bar.fill.fore_color.rgb = YELLOW; bar.line.fill.background()
    add_slide_number(s, "10")
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(6), Inches(0.4), "NEW CHAPTER  •  STUDENT AMBASSADOR PROGRAM", 10, True, YELLOW, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(6.5), Inches(1.6), "Carry the AI Yatra to your campus.", 36, True, WHITE, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(2.4), Inches(6.2), Inches(0.9), "Ambassadors rally their batch to Saturday meetups, host campus mini-sessions, and turn dense AI into hands-on learning — free, always.", 13, False, WHITE, PP_ALIGN.LEFT)
    duties = ["📣  Campus outreach — posters, clubs, WhatsApp", "👥  Bring your batch — help first-timers settle in", "📅  Host mini-sessions — paper-readings & build nights", "📸  Document — photos, clips, field notes", "📖  Learn fast — first access + prep calls", "🧭  Shape roadmap — student feedback → meetup arc"]
    add_bullet_list(s, Inches(0.6), Inches(3.5), Inches(6.2), Inches(2.6), ["•  " + d for d in duties], 11, WHITE)
    # right card
    card = add_rounded_rect(s, Inches(7.4), Inches(0.7), Inches(5.3), Inches(5.6), WHITE)
    add_textbox(s, Inches(7.9), Inches(1.0), Inches(4.4), Inches(0.4), "WHAT AMBASSADORS GET", 10, True, DEEP, PP_ALIGN.LEFT)
    perks = [("Official recognition", "Signed certificate + public profile after first term."), ("Letters that count", "Rec letters for internships & higher studies — earned by shipped work."), ("Direct mentor access", "Monthly office hours with organizers."), ("Speaker lane", "Priority demo / lightning-talk / lead slots.")]
    y = 1.5
    for t, b in perks:
        add_textbox(s, Inches(7.9), Inches(y), Inches(4.4), Inches(0.35), "✓  " + t, 11.5, True, INK, PP_ALIGN.LEFT)
        add_textbox(s, Inches(7.9), Inches(y+0.35), Inches(4.4), Inches(0.45), b, 10, False, SLATE, PP_ALIGN.LEFT)
        y += 0.85
    add_textbox(s, Inches(7.9), Inches(5.15), Inches(4.4), Inches(0.7), f"Apply: {EMAIL}\nSubject: AIYatra Student Ambassador — Application", 10, True, NAVY, PP_ALIGN.LEFT)
    add_footer(s)

    # ---- 11. Join us ----
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_slide_number(s, "11")
    add_textbox(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.8), "Your AI journey starts with a single RSVP.", 36, True, INK, PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.5), Inches(1.4), Inches(8.3), Inches(0.5), f"Join 2,957 learners democratizing AI — one meetup, one project, one breakthrough at a time.", 13, False, SLATE, PP_ALIGN.CENTER)
    contacts = [
        ("WEBSITE", WEBSITE, "Gatherings · Blog · Ambassadors"),
        ("LINKEDIN", "linkedin.com/company/aiyatra", "Follow milestones & sessions"),
        ("EMAIL", EMAIL, "Ambassadors + questions"),
        ("MEETUP", "meetup.com/aiyatra", "RSVP: Sep 5 Agentic AI"),
    ]
    for i, (k, v, d) in enumerate(contacts):
        x = Inches(0.6 + i*3.05)
        card = add_rounded_rect(s, x, Inches(2.2), Inches(2.85), Inches(2.8), WHITE, MUTED, 1)
        hdr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(2.85), Inches(0.55))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = NAVY; hdr.line.fill.background()
        hdr.text_frame.word_wrap = True; hdr.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = hdr.text_frame.paragraphs[0]; p.text = k; p.font.size = Pt(10); p.font.bold=True; p.font.color.rgb=YELLOW; p.font.name="Calibri"; p.alignment=PP_ALIGN.CENTER
        add_textbox(s, x+Inches(0.2), Inches(2.95), Inches(2.45), Inches(0.7), v, 11, True, INK, PP_ALIGN.CENTER)
        add_textbox(s, x+Inches(0.2), Inches(3.65), Inches(2.45), Inches(0.6), d, 10, False, SLATE, PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(5.3), Inches(12.13), Inches(0.6), "Free forever · No prerequisites  •  Saturdays, mornings IST  •  LSEG, International Tech Park, Madhapur, Hyderabad  •  © 2026 AIYatra · Research. Build. Transform.", 10, True, MUTED, PP_ALIGN.CENTER)
    add_footer(s)

    path = os.path.join(OUT_DIR, "AIYatra-Journey-From-Start-Till-Now.pptx")
    prs.save(path)
    print("Saved:", path)
    return path

# ============================================================
# DECK 2 — Ambassador Program
# ============================================================
def build_deck2():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Pt(10))
    bar.fill.solid(); bar.fill.fore_color.rgb = YELLOW; bar.line.fill.background()
    add_textbox(s, Inches(0.8), Inches(0.7), Inches(7.5), Inches(0.5), "AIYATRA  •  NEW CHAPTER  •  OPEN TO EVERY COLLEGE", 11, True, YELLOW, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.8), Inches(1.15), Inches(7.5), Inches(1.6), "Student Ambassador Program", 52, True, WHITE, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.8), Inches(2.7), Inches(6.8), Inches(0.6), "Carry the AI Yatra to your campus.", 26, True, YELLOW, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.8), Inches(3.4), Inches(6.8), Inches(1.0), "Be the face of the movement in your college — rally classmates to Saturday meetups, host campus mini-sessions, and turn dense AI into hands-on learning.", 14, False, WHITE, PP_ALIGN.LEFT)
    card = add_rounded_rect(s, Inches(8.4), Inches(0.9), Inches(4.1), Inches(5.0), WHITE)
    add_textbox(s, Inches(8.8), Inches(1.2), Inches(3.3), Inches(0.4), "APPLY IN ONE EMAIL", 10, True, DEEP, PP_ALIGN.LEFT)
    add_textbox(s, Inches(8.8), Inches(1.6), Inches(3.3), Inches(0.6), EMAIL, 13, True, NAVY, PP_ALIGN.LEFT)
    add_textbox(s, Inches(8.8), Inches(2.25), Inches(3.3), Inches(1.0), "Subject: AIYatra Student Ambassador — Application\n\nInclude: Name · College · Year/Branch · Why you (1 para)", 10.5, False, SLATE, PP_ALIGN.LEFT)
    add_textbox(s, Inches(8.8), Inches(3.5), Inches(3.3), Inches(0.8), "✓ Any year, any branch\n✓ ~3 hrs/week\n✓ Free, always", 11, True, INK, PP_ALIGN.LEFT)
    add_textbox(s, Inches(8.8), Inches(4.45), Inches(3.3), Inches(1.0), f"{WEBSITE}/ambassadors\nlinkedin.com/company/aiyatra\nmeetup.com/aiyatra", 10, False, SLATE, PP_ALIGN.LEFT)
    add_footer(s)

    # 2 Why
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_top_kicker(s, "Why ambassadors  •  why now")
    add_slide_number(s, "02")
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.7), "Most AI events feel out of reach. Ambassadors fix that.", 28, True, INK, PP_ALIGN.LEFT)
    whys = [
        ("THE GAP", "Students hear about AIYatra late, come alone once, and drift away. Campuses need a familiar face.", YELLOW),
        ("THE BRIDGE", "You are that bridge — posters, classroom shout-outs, WhatsApp groups, tech clubs.", TEAL),
        ("THE PAYOFF", "More first-timers who stay, stronger campus study groups, and a community that grows one batch at a time.", CORAL),
    ]
    for i, (t, b, accent) in enumerate(whys):
        x = Inches(0.6 + i*4.05)
        card = add_rounded_rect(s, x, Inches(1.9), Inches(3.75), Inches(2.9), WHITE, MUTED, 1)
        stripe = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Inches(0.3), Inches(2.15), Inches(0.7), Pt(10))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = accent; stripe.line.fill.background()
        add_textbox(s, x+Inches(0.3), Inches(2.5), Inches(3.15), Inches(0.4), t, 11, True, INK, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.3), Inches(2.95), Inches(3.15), Inches(1.5), b, 12, False, SLATE, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(5.1), Inches(12.13), Inches(0.9), "Context: 2,957 members · 11 Saturday sessions at LSEG Madhapur · 4.6★ (90 ratings) — and students consistently say AIYatra is the first AI room that felt welcoming and deep.", 11.5, False, SLATE, PP_ALIGN.LEFT)
    add_footer(s)

    # 3 Six duties
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_top_kicker(s, "The mandate  •  six duties, one movement")
    add_slide_number(s, "03")
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6), "Six duties. One semester. Renew with honors.", 28, True, INK, PP_ALIGN.LEFT)
    duties = [
        ("1 · CARRY TO CAMPUS", "Posters, shout-outs, WhatsApp & clubs. Bridge AIYatra ↔ campus.", YELLOW),
        ("2 · BRING YOUR BATCH", "Rally classmates to Saturdays; help first-timers from RSVP to notes.", TEAL),
        ("3 · HOST MINI-SESSIONS", "Paper-readings, build nights, revision circles with our kits.", NAVY),
        ("4 · DOCUMENT", "Photos, clips, field notes — grow the story online.", CORAL),
        ("5 · LEARN FAST", "First access to material + prep calls: transformers, PyTorch, agents.", LIGHT_GREY),
        ("6 · FEED BACK", "Tell us what students want — topics, timings shape the arc.", LIGHT_GREY),
    ]
    for i, (t, b, fill) in enumerate(duties):
        col = i % 3; row = i // 3
        x = Inches(0.6 + col*4.05); y = 1.75 + row*2.35
        is_navy = (fill == NAVY)
        tc = WHITE if is_navy else INK
        bc = WHITE if is_navy else SLATE
        card = add_rounded_rect(s, x, Inches(y), Inches(3.75), Inches(2.1), fill, MUTED if not is_navy else NAVY, 1)
        add_textbox(s, x+Inches(0.3), Inches(y+0.2), Inches(3.15), Inches(0.5), t, 11, True, CORAL if (not is_navy and fill==LIGHT_GREY) else (YELLOW if is_navy else DEEP), PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.3), Inches(y+0.75), Inches(3.15), Inches(1.1), b, 11.5, False, bc, PP_ALIGN.LEFT)
    add_footer(s)

    # 4 Perks
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    add_slide_number(s, "04")
    add_textbox(s, Inches(0.6), Inches(0.4), Inches(6), Inches(0.4), "WHAT YOU GET  •  PROOF OF WORK, NOT SWAG", 10, True, YELLOW, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7), "Ambassadorship pays in proof of work.", 32, True, WHITE, PP_ALIGN.LEFT)
    perks = [
        ("Official recognition", "Signed Ambassador certificate + public profile on aiyATRA.io/ambassadors after your first term."),
        ("Letters that count", "Recommendation letters for internships & higher studies — earned through shipped work, not attendance."),
        ("Direct mentor access", "Monthly office hours with organizers. Projects, careers, wild paper ideas — bring them."),
        ("Speaker lane", "Priority slots to demo, lightning-talk, and eventually lead Saturday sessions."),
    ]
    for i, (t, b) in enumerate(perks):
        x = Inches(0.6 + (i%2)*6.37)
        y = 1.8 + (i//2)*2.2
        card = add_rounded_rect(s, x, Inches(y), Inches(6.16), Inches(1.9), WHITE)
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Inches(0.3), Inches(y+0.3), Inches(0.9), Inches(0.45))
        badge.fill.solid(); badge.fill.fore_color.rgb = YELLOW; badge.line.fill.background()
        badge.text_frame.word_wrap=True; badge.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=badge.text_frame.paragraphs[0]; p.text=f"0{i+1}"; p.font.size=Pt(11); p.font.bold=True; p.font.color.rgb=INK; p.font.name="Calibri"; p.alignment=PP_ALIGN.CENTER
        add_textbox(s, x+Inches(1.4), Inches(y+0.25), Inches(4.4), Inches(0.5), t, 15, True, INK, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(1.4), Inches(y+0.75), Inches(4.4), Inches(0.9), b, 11, False, SLATE, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(6.35), Inches(12.13), Inches(0.3), "No stipends, no swag-bribes — everything compounds into skills, relationships, and verifiable receipts.", 10, False, WHITE, PP_ALIGN.CENTER)
    add_footer(s)

    # 5 Who fits
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_top_kicker(s, "Who fits  •  built for students who show up")
    add_slide_number(s, "05")
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(7), Inches(0.7), "Curiosity beats CGPA here.", 30, True, INK, PP_ALIGN.LEFT)
    fits = [
        "Enrolled in any college / university — any year, any branch.",
        "Attended ≥1 AIYatra meetup (or will attend the next one) — know the room you invite people into.",
        "Can give ~3 hours/week: one meetup + campus outreach + short weekly check-in.",
        "Comfortable talking to people — classrooms, clubs, group chats — and following through.",
    ]
    box = add_textbox(s, Inches(0.6), Inches(1.75), Inches(7.0), Inches(0.4), "", 12, False, INK, PP_ALIGN.LEFT)
    tf = box.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; p.text="✓  "+fits[0]; p.font.size=Pt(12); p.font.color.rgb=INK; p.font.name="Calibri"; p.space_after=Pt(10)
    for f in fits[1:]:
        add_para(tf, "✓  "+f, 12, False, INK, PP_ALIGN.LEFT, Pt(10))
    card = add_rounded_rect(s, Inches(0.6), Inches(4.3), Inches(7.0), Inches(1.9), YELLOW, INK, 1.5)
    add_textbox(s, Inches(1.0), Inches(4.5), Inches(6.2), Inches(0.4), "♥  Missing one? Apply anyway.", 15, True, INK, PP_ALIGN.LEFT)
    add_textbox(s, Inches(1.0), Inches(4.95), Inches(6.2), Inches(1.0), "Most ambassadors started as first-timers in the back row. No speaking experience or big following needed — write to us, we'll figure it out together.", 11, False, INK, PP_ALIGN.LEFT)
    # right: time box
    card2 = add_rounded_rect(s, Inches(8.1), Inches(1.75), Inches(4.0), Inches(4.45), WHITE, MUTED, 1)
    add_textbox(s, Inches(8.5), Inches(2.0), Inches(3.2), Inches(0.4), "A WEEK IN THE LIFE  (~3 HRS)", 10, True, DEEP, PP_ALIGN.LEFT)
    weeks = [("SAT  —  Meetup", "Attend + bring classmates, help newcomers."), ("MIDWEEK  —  Outreach", "Posters, club sync, WhatsApp nudges."), ("CHECK-IN  —  20 min", "Async update: wins, blocks, next drive.")]
    y=2.6
    for t,b in weeks:
        add_textbox(s, Inches(8.5), Inches(y), Inches(3.2), Inches(0.35), t, 11, True, INK, PP_ALIGN.LEFT)
        add_textbox(s, Inches(8.5), Inches(y+0.35), Inches(3.2), Inches(0.5), b, 10, False, SLATE, PP_ALIGN.LEFT)
        y+=0.95
    add_textbox(s, Inches(8.5), Inches(5.45), Inches(3.2), Inches(0.4), "Term = one semester", 10, True, MUTED, PP_ALIGN.LEFT)
    add_footer(s)

    # 6 How to apply
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_top_kicker(s, "How to apply  •  three steps to the crest")
    add_slide_number(s, "06")
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.6), "Three steps to the crest.", 30, True, INK, PP_ALIGN.LEFT)
    steps = [
        ("01 · SAY HELLO", f"Email {EMAIL} — name, college, year, + 1 para on why you want to carry the yatra.", YELLOW),
        ("02 · MEET THE CREW", "Short intro call with an organizer — campus scene + your first 30-day plan.", TEAL),
        ("03 · RUN YOUR FIRST DRIVE", "Bring 5 classmates to a Saturday + host 1 campus huddle. Do that → crest is yours.", NAVY),
    ]
    for i,(t,b,fill) in enumerate(steps):
        x = Inches(0.6+i*4.05)
        is_navy=(fill==NAVY)
        card = add_rounded_rect(s, x, Inches(1.9), Inches(3.75), Inches(2.9), fill, MUTED if not is_navy else NAVY, 1)
        tc = WHITE if is_navy else INK; bc = WHITE if is_navy else SLATE
        add_textbox(s, x+Inches(0.3), Inches(2.15), Inches(3.15), Inches(0.6), t, 12, True, (YELLOW if is_navy else DEEP), PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.3), Inches(2.85), Inches(3.15), Inches(1.4), b, 12, False, bc, PP_ALIGN.LEFT)
    # email template strip
    strip = add_rounded_rect(s, Inches(0.6), Inches(5.1), Inches(12.13), Inches(1.1), CREAM, MUTED, 1)
    add_textbox(s, Inches(1.0), Inches(5.25), Inches(11.3), Inches(0.8),
                f"EMAIL TEMPLATE  •  To: {EMAIL}  •  Subject: AIYatra Student Ambassador — Application  •  Body: Name / College / Year-Branch / Why I want to be an ambassador / Meetups attended",
                10.5, True, SLATE, PP_ALIGN.LEFT)
    add_footer(s)

    # 7 Campus playbook
    s = prs.slides.add_slide(blank)
    set_bg(s, CREAM)
    add_top_kicker(s, "Campus playbook  •  what mini-sessions look like")
    add_slide_number(s, "07")
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(12), Inches(0.7), "Same rhythm, your turf: research · build · transform.", 28, True, INK, PP_ALIGN.LEFT)
    plays = [
        ("PAPER-READINGS", "45–60 min. 1 paper section + discussion. Zero lectures. Kit provided.", "e.g. Attention Is All You Need, §3"),
        ("BUILD NIGHTS", "2 hrs. Laptops open — tensors, small agent loops, verified PRs.", "e.g. Rebuild nn.Linear together"),
        ("REVISION CIRCLES", "Post-Saturday recap for classmates who missed the meetup.", "e.g. DeepSeek-V3 in 30 min"),
        ("DEMO HOURS", "Lightning talks by batchmates. Ambassadors get first slots.", "e.g. My first MCP extension"),
    ]
    for i,(t,b,e) in enumerate(plays):
        x = Inches(0.6+i*3.05)
        card = add_rounded_rect(s, x, Inches(1.9), Inches(2.85), Inches(3.2), WHITE, MUTED, 1)
        add_textbox(s, x+Inches(0.3), Inches(2.15), Inches(2.25), Inches(0.4), t, 10, True, DEEP, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.3), Inches(2.6), Inches(2.25), Inches(1.2), b, 11, False, SLATE, PP_ALIGN.LEFT)
        add_textbox(s, x+Inches(0.3), Inches(3.85), Inches(2.25), Inches(0.6), e, 9.5, True, MUTED, PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.6), Inches(5.4), Inches(12.13), Inches(0.6), "We provide curriculum kits + prep calls. You provide the room, the friends, and the follow-through.", 11, True, SLATE, PP_ALIGN.LEFT)
    add_footer(s)

    # 8 Impact + FAQ
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    add_top_kicker(s, "Impact + honest answers")
    add_slide_number(s, "08")
    add_textbox(s, Inches(0.6), Inches(0.9), Inches(7), Inches(0.6), "Good for you. Great for campus.", 26, True, INK, PP_ALIGN.LEFT)
    impacts = ["Leadership you can prove — drives, sessions, docs.", "Mentorship + speaker lane — office hours to stage.", "Network — organizers, builders, founders, study groups.", "Letters that say something — because you shipped."]
    add_bullet_list(s, Inches(0.6), Inches(1.65), Inches(6.5), Inches(2.2), ["•  "+x for x in impacts], 12, SLATE)
    # FAQ card
    card = add_rounded_rect(s, Inches(7.5), Inches(1.6), Inches(5.2), Inches(4.5), NAVY)
    add_textbox(s, Inches(7.9), Inches(1.85), Inches(4.4), Inches(0.4), "HONEST ANSWERS", 10, True, YELLOW, PP_ALIGN.LEFT)
    faqs = [("Stipend?", "No. Proof-of-work: skills, letters, stage."), ("Who can apply?", "Any college, any year, any branch."), ("Time?", "~3 hrs/week for one semester."), ("Experience needed?", "None. Curiosity + follow-through."), ("Cost?", "Free, always.")]
    y=2.3
    for q,a in faqs:
        add_textbox(s, Inches(7.9), Inches(y), Inches(4.4), Inches(0.35), q+"  →  "+a, 10.5, False, WHITE, PP_ALIGN.LEFT)
        y+=0.45
    add_textbox(s, Inches(7.9), Inches(5.3), Inches(4.4), Inches(0.5), "Questions? Ask us anything — same email.", 10, True, YELLOW, PP_ALIGN.LEFT)
    add_footer(s)

    # 9 CTA
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Pt(10))
    bar.fill.solid(); bar.fill.fore_color.rgb = YELLOW; bar.line.fill.background()
    add_slide_number(s, "09")
    add_textbox(s, Inches(0.6), Inches(0.7), Inches(12.13), Inches(0.5), "ONE CAMPUS AT A TIME", 11, True, YELLOW, PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(1.1), Inches(12.13), Inches(1.2), "Your campus is the next AI Yatra.", 44, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.8), Inches(2.3), Inches(7.7), Inches(0.6), "Write today — or come experience a Saturday first, then decide.", 14, False, WHITE, PP_ALIGN.CENTER)
    contacts = [
        ("EMAIL", EMAIL, "Application + questions"),
        ("WEB", f"{WEBSITE}/ambassadors", "Program details"),
        ("LINKEDIN", "linkedin.com/company/aiyatra", "Follow + DM"),
        ("MEETUP", "meetup.com/aiyatra", "Attend first: Sep 5"),
    ]
    for i,(k,v,d) in enumerate(contacts):
        x = Inches(0.6+i*3.05)
        card = add_rounded_rect(s, x, Inches(3.2), Inches(2.85), Inches(2.3), WHITE)
        add_textbox(s, x+Inches(0.2), Inches(3.4), Inches(2.45), Inches(0.3), k, 9, True, DEEP, PP_ALIGN.CENTER)
        add_textbox(s, x+Inches(0.2), Inches(3.7), Inches(2.45), Inches(0.7), v, 10, True, INK, PP_ALIGN.CENTER)
        add_textbox(s, x+Inches(0.2), Inches(4.4), Inches(2.45), Inches(0.4), d, 9, False, MUTED, PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(5.85), Inches(12.13), Inches(0.4), "© 2026 AIYatra  •  Research. Build. Transform.  •  Free, always  •  Hyderabad + every campus", 10, True, YELLOW, PP_ALIGN.CENTER)
    add_footer(s)

    path = os.path.join(OUT_DIR, "AIYatra-Student-Ambassador-Program.pptx")
    prs.save(path)
    print("Saved:", path)
    return path

if __name__ == "__main__":
    p1 = build_deck1()
    p2 = build_deck2()
    print("Done.")
