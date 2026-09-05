#!/usr/bin/env python3
"""AIYatra decks v2 — faithful port of aiyatra.io theme.

Theme source: apps/web/src/index.css + tailwind.config.js + HomePage/AmbassadorsPage/SiteChrome
- paper #F5F1E5, paper-soft #F9F7F1, ink #272520, ink-soft #68635A
- tones: blue #CEE1EE, blue-deep #477A9E, green #D4E7D0, yellow #F5E4AD, coral #F2D6CA, violet #E0D8EE
- fonts: Caveat (font-hand, headings) + DM Sans (body) — same names as website
- cards: rounded, 2px ink border, hard offset shadow (5px 5px 0 ink), no gradients
- header: paper-soft bar, ink bottom border, logo + nav + blue-deep Join button
- footer: ink bar, 5px yellow top border, paper text
- copy: verbatim from website components
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUT_DIR = "/Users/moinuddin/Documents/AIYatra/aiyatra.io-website/AIYatra.IO/presentations"
LOGO = "/Users/moinuddin/Documents/AIYatra/aiyatra.io-website/AIYatra.IO/apps/web/public/aiyatra-mark.png"
CREST = "/Users/moinuddin/Documents/AIYatra/aiyatra.io-website/AIYatra.IO/apps/web/public/ambassador-crest.jpg"

os.makedirs(OUT_DIR, exist_ok=True)

# ---- exact website palette (HSL -> RGB) ----
PAPER      = RGBColor(0xF5, 0xF1, 0xE5)
PAPER_SOFT = RGBColor(0xF9, 0xF7, 0xF1)
INK        = RGBColor(0x27, 0x25, 0x20)
INK_SOFT   = RGBColor(0x68, 0x63, 0x5A)
T_BLUE     = RGBColor(0xCE, 0xE1, 0xEE)
T_BLUE_D   = RGBColor(0x47, 0x7A, 0x9E)
T_GREEN    = RGBColor(0xD4, 0xE7, 0xD0)
T_YELLOW   = RGBColor(0xF5, 0xE4, 0xAD)
T_CORAL    = RGBColor(0xF2, 0xD6, 0xCA)
T_VIOLET   = RGBColor(0xE0, 0xD8, 0xEE)
# section tints (tone @30-40% over paper, as in site: bg-tone-blue/30, bg-tone-green/40)
EVENTS_BG  = RGBColor(0xE9, 0xEC, 0xE7)  # tone-blue/30 over paper
ABOUT_BG   = RGBColor(0xE7, 0xED, 0xDC)  # tone-green/40 over paper
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

HEAD = "Caveat"
BODY = "DM Sans"

WEBSITE = "https://aiyatra.io"
WEBSITE_SHORT = "aiyatra.io"
LINKEDIN = "https://www.linkedin.com/company/aiyatra/"
LINKEDIN_SHORT = "linkedin.com/company/aiyatra"
EMAIL = "global.aiyatra@gmail.com"
MEETUP = "https://www.meetup.com/aiyatra/"
MEETUP_SHORT = "meetup.com/aiyatra"
GITHUB_SHORT = "github.com/AI-Yatra"

FOOTER_LINE = f"{WEBSITE_SHORT}   •   {LINKEDIN_SHORT}   •   {EMAIL}   •   {MEETUP_SHORT}"

SW, SH = Inches(13.33), Inches(7.5)
GX = Inches(0.55)  # gutter matching .wrap
GW = Inches(12.23)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rounded(slide, l, t, w, h, fill, line=None, lw=1.5, shadow=True, radius=0.12):
    if shadow:
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l + Pt(5), t + Pt(5), w, h)
        sh.fill.solid(); sh.fill.fore_color.rgb = INK
        sh.line.fill.background()
        try: sh.adjustments[0] = radius
        except: pass
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    try: s.adjustments[0] = radius
    except: pass
    return s

def tb(slide, l, t, w, h, text, size=11, bold=False, color=INK, font=BODY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font; p.alignment = align
    return box

def para(tf, text, size=11, bold=False, color=INK, font=BODY, align=PP_ALIGN.LEFT, after=Pt(3)):
    p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font; p.alignment = align
    p.space_after = after; p.space_before = Pt(1)
    return p

def pill(slide, l, t, w, h, text, fill=T_YELLOW):
    _rounded(slide, l, t, w, h, fill, INK, 1.25, shadow=True, radius=0.5)
    b = slide.shapes.add_textbox(l, t, w, h)
    b.text_frame.word_wrap = True
    b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = b.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(9); p.font.bold = True
    p.font.color.rgb = INK; p.font.name = BODY; p.alignment = PP_ALIGN.CENTER
    return b

def btn(slide, l, t, w, h, text, fill=T_BLUE_D, tcolor=PAPER_SOFT):
    _rounded(slide, l, t, w, h, fill, INK, 1.5, shadow=True, radius=0.18)
    b = slide.shapes.add_textbox(l, t, w, h)
    b.text_frame.word_wrap = True
    b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = b.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(10.5); p.font.bold = True
    p.font.color.rgb = tcolor; p.font.name = BODY; p.alignment = PP_ALIGN.CENTER
    return b

def header(slide, dark=False):
    # top bar like SiteChrome Header: paper bar, ink bottom border
    bar_bg = INK if dark else PAPER_SOFT
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SW, Inches(0.72))
    bar.fill.solid(); bar.fill.fore_color.rgb = bar_bg; bar.line.fill.background()
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.72), SW, Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = T_YELLOW if dark else INK; ln.line.fill.background()
    tc = PAPER_SOFT if dark else INK
    sc = T_YELLOW if dark else INK_SOFT
    # logo
    try:
        slide.shapes.add_picture(LOGO, GX, Inches(0.08), Inches(0.56), Inches(0.51))
    except: pass
    tb(slide, GX + Inches(0.64), Inches(0.10), Inches(1.6), Inches(0.32), "AI Yatra", 13, True, tc, BODY)
    tb(slide, GX + Inches(0.64), Inches(0.34), Inches(1.6), Inches(0.22), "RESEARCH · BUILD · TRANSFORM", 6.5, True, sc, BODY)
    # nav (same labels as NAV_LINKS)
    nav = "Events     Community     Approach     Testimonials     About     Ambassadors     Blog"
    nbox = tb(slide, Inches(3.6), Inches(0.14), Inches(6.4), Inches(0.44), nav, 8, False, sc, BODY, PP_ALIGN.CENTER)
    nbox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # join button
    btn(slide, Inches(10.6), Inches(0.12), Inches(2.18), Inches(0.48), "Join on Meetup ↗")

def footer(slide, dark_section=False):
    # website footer replica: ink bar, 5px yellow top border
    y0 = SH - Inches(0.5)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), y0, SW, Pt(5))
    top.fill.solid(); top.fill.fore_color.rgb = T_YELLOW; top.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), y0 + Pt(5), SW, Inches(0.5) - Pt(5))
    bar.fill.solid(); bar.fill.fore_color.rgb = INK; bar.line.fill.background()
    tb(slide, Inches(0.3), y0 + Pt(8), Inches(10.5), Inches(0.3), FOOTER_LINE, 7.5, False, PAPER_SOFT, BODY, PP_ALIGN.LEFT)
    tb(slide, Inches(11.0), y0 + Pt(8), Inches(1.9), Inches(0.3), "© 2026 AIYatra", 7.5, False, PAPER_SOFT, BODY, PP_ALIGN.RIGHT)

def kicker_line(slide, y, text, color=T_BLUE_D):
    tb(slide, GX, y, GW, Inches(0.3), text, 9, True, color, BODY)

def h_hand(slide, y, text, size=34, color=INK, h=Inches(0.7), align=PP_ALIGN.LEFT):
    return tb(slide, GX, y, GW, h, text, size, True, color, HEAD, align)

# ================= DECK 1 =================
def deck1():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    blank = prs.slide_layouts[6]

    # -- 1. Hero (paper) : exact Hero copy --
    s = prs.slides.add_slide(blank); bg(s, PAPER); header(s)
    pill(s, GX, Inches(1.0), Inches(2.9), Inches(0.38), "✦  Democratizing AI Learning")
    tb(s, GX, Inches(1.48), Inches(7.0), Inches(1.5),
       "Research. Build. Transform.", 46, True, INK, HEAD)
    b = tb(s, GX, Inches(2.75), Inches(6.9), Inches(0.9),
       "AIYatra is Hyderabad's open AI community — a place where beginners and experts share knowledge, ship real projects, and master machine learning together. No paywalls, no prerequisites. Just curiosity and a laptop.",
       11.5, False, INK_SOFT, BODY)
    btn(s, GX, Inches(3.85), Inches(2.6), Inches(0.52), "RSVP to the next meetup  →")
    _rounded(s, GX + Inches(2.8), Inches(3.85), Inches(2.5), Inches(0.52), PAPER_SOFT, INK, 1.5, True, 0.18)
    tb(s, GX + Inches(2.8), Inches(3.85), Inches(2.5), Inches(0.52), "Explore past events", 10.5, True, INK, BODY, PP_ALIGN.CENTER).text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # stats row (CountUp values)
    for i, (v, l) in enumerate([("2,957", "Community members"), ("14", "Events hosted"), ("4.6★", "From 90 ratings")]):
        x = GX + Inches(i * 2.1)
        tb(s, x, Inches(4.65), Inches(1.9), Inches(0.5), v, 30, True, INK, HEAD)
        tb(s, x, Inches(5.15), Inches(1.9), Inches(0.3), l.upper(), 8, True, INK_SOFT, BODY)
    # right HeroCard (exact)
    cx, cy, cw, ch = Inches(7.9), Inches(1.0), Inches(4.85), Inches(5.15)
    _rounded(s, cx, cy, cw, ch, PAPER_SOFT, INK, 1.75, True, 0.08)
    pill(s, cx + Inches(0.3), cy + Inches(0.25), Inches(1.6), Inches(0.34), "🎟  Next meetup")
    tb(s, cx + Inches(2.2), Inches(0.25) + cy, Inches(2.3), Inches(0.34), "● RSVP open", 8.5, True, T_BLUE_D, BODY, PP_ALIGN.RIGHT)
    tb(s, cx + Inches(0.3), cy + Inches(0.7), Inches(4.25), Inches(0.7), "Harnessing Agentic AI", 30, True, INK, HEAD)
    tb(s, cx + Inches(0.3), cy + Inches(1.4), Inches(4.25), Inches(0.7),
       "Build the loop, the tools and the guardrails — and make your own coding agent that reads, plans, edits, tests and repairs its way to a verified pull request. 100% offline, zero API keys.",
       9.5, False, INK_SOFT, BODY)
    _rounded(s, cx + Inches(0.3), cy + Inches(2.25), Inches(2.0), Inches(0.95), T_BLUE, INK, 1.25, False, 0.15)
    tb(s, cx + Inches(0.45), cy + Inches(2.3), Inches(1.8), Inches(0.25), "WHEN", 7.5, True, INK_SOFT, BODY)
    tb(s, cx + Inches(0.45), cy + Inches(2.55), Inches(1.8), Inches(0.5), "Sat, Sep 5 · 9:00 AM IST", 10, True, INK, BODY)
    _rounded(s, cx + Inches(2.55), cy + Inches(2.25), Inches(2.0), Inches(0.95), T_GREEN, INK, 1.25, False, 0.15)
    tb(s, cx + Inches(2.7), cy + Inches(2.3), Inches(1.8), Inches(0.25), "ATTENDING", 7.5, True, INK_SOFT, BODY)
    tb(s, cx + Inches(2.7), cy + Inches(2.55), Inches(1.8), Inches(0.5), "319 learners & counting", 10, True, INK, BODY)
    for i, t in enumerate(["The agent loop — the core inside Claude Code & Codex", "Tools, planning & self-repair in pure Python", "Guardrails & sandboxing — 100% offline"]):
        tb(s, cx + Inches(0.55), cy + Inches(3.35 + i * 0.32), Inches(4.0), Inches(0.32), f"{i+1}   {t}", 8.5, False, INK_SOFT, BODY)
    btn(s, cx + Inches(0.3), cy + Inches(4.4), Inches(4.25), Inches(0.5), "RSVP now  →")
    footer(s)

    # -- 2. Events (events bg = tone-blue/30) --
    s = prs.slides.add_slide(blank); bg(s, EVENTS_BG); header(s)
    kicker_line(s, Inches(1.0), "●  EVENTS  ·  PAST SESSIONS")
    h_hand(s, Inches(1.28), "Hands-on Saturdays, archived on Meetup.", 32)
    cards = [
        ("Past event", "PyTorch Foundations", "Tensors, autograd and the training loop — the foundations every AI engineer stands on. Rebuilt nn.Linear from scratch.", "Sat, Aug 22 · 10:00 AM IST · 146 attended", T_BLUE),
        ("Past event", "Linear Algebra → Transformers", "Vectors to attention: embeddings, Q/K/V projections, LoRA and SVD — the math beneath the models, by hand.", "Sat, Aug 15 · 10:00 AM IST · 207 attended", T_YELLOW),
        ("Past event", "DeepSeek v3 from Scratch", "Multi-head latent attention, MoE load balancing and RoPE scaling — the DeepSeek-V3 paper, live-coded in PyTorch.", "Sat, Aug 8 · 10:00 AM IST · 286 attended", T_VIOLET),
    ]
    for i, (tag, title, blurb, meta, tone) in enumerate(cards):
        x = GX + Inches(i * 4.15)
        _rounded(s, x, Inches(2.2), Inches(3.93), Inches(2.9), PAPER_SOFT, INK, 1.5, True, 0.08)
        pill(s, x + Inches(0.25), Inches(2.4), Inches(1.5), Inches(0.32), f"◷  {tag}", tone)
        tb(s, x + Inches(0.25), Inches(2.85), Inches(3.43), Inches(0.6), title, 21, True, INK, HEAD)
        tb(s, x + Inches(0.25), Inches(3.5), Inches(3.43), Inches(0.75), blurb, 9.5, False, INK_SOFT, BODY)
        tb(s, x + Inches(0.25), Inches(4.35), Inches(3.43), Inches(0.4), meta, 8.5, True, INK, BODY)
    # archive band (exact copy)
    _rounded(s, GX, Inches(5.35), GW, Inches(1.1), PAPER_SOFT, INK, 1.5, True, 0.1)
    tb(s, GX + Inches(0.4), Inches(5.5), Inches(7.5), Inches(0.35), "The full archive lives on Meetup.", 20, True, INK, HEAD)
    tb(s, GX + Inches(0.4), Inches(5.95), Inches(7.5), Inches(0.3), "Photos, discussions, ratings and RSVPs — straight from the source.", 10, False, INK_SOFT, BODY)
    btn(s, GX + Inches(9.0), Inches(5.65), Inches(2.8), Inches(0.5), "View all on Meetup ↗", PAPER_SOFT, INK)
    footer(s)

    # -- 3. Community (paper-soft, border-y ink) --
    s = prs.slides.add_slide(blank); bg(s, PAPER_SOFT); header(s)
    kicker_line(s, Inches(1.0), "◉  COMMUNITY")
    h_hand(s, Inches(1.28), "Real rooms. Real laptops. Real people.", 32)
    tb(s, GX, Inches(1.95), Inches(8.5), Inches(0.35),
       "10 hands-on sessions since Hyderabad — the numbers tell the story of rooms at LSEG, International Tech Park, Madhapur where 2,957 members learn together.",
       10.5, False, INK_SOFT, BODY)
    # col 1 biggest rooms (tone-yellow)
    _rounded(s, GX, Inches(2.5), Inches(3.93), Inches(3.4), T_YELLOW, INK, 1.5, True, 0.08)
    tb(s, GX + Inches(0.3), Inches(2.65), Inches(3.3), Inches(0.3), "◉  BIGGEST ROOMS SO FAR", 8, True, INK_SOFT, BODY)
    for i, (t, m) in enumerate([("Harnessing Agentic AI — 319", "#1 by turnout · Sat, Sep 5"), ("DeepSeek v3 from Scratch — 286", "#2 by turnout · Sat, Aug 8"), ("Linear Algebra → Transformers — 207", "#3 by turnout · Sat, Aug 15")]):
        _rounded(s, GX + Inches(0.3), Inches(3.05 + i * 0.75), Inches(3.33), Inches(0.62), PAPER_SOFT, INK, 1.25, False, 0.12)
        tb(s, GX + Inches(0.5), Inches(3.1 + i * 0.75), Inches(3.0), Inches(0.3), t, 9, True, INK, BODY)
        tb(s, GX + Inches(0.5), Inches(3.32 + i * 0.75), Inches(3.0), Inches(0.25), m, 8, False, INK_SOFT, BODY)
    # col 2 learning tracks (paper)
    _rounded(s, GX + Inches(4.15), Inches(2.5), Inches(3.93), Inches(3.4), PAPER, INK, 1.5, True, 0.08)
    tb(s, GX + Inches(4.45), Inches(2.65), Inches(3.3), Inches(0.3), "▤  LEARNING TRACKS", 8, True, INK_SOFT, BODY)
    for i, t in enumerate(["Foundations. Transformer paper → code, linear algebra → attention, PyTorch tensors → training loop.",
                            "Agents. Goose end-to-end demo, DeepSeek-V3 architecture, and the Sep 5 coding-agent build.",
                            "Reading room. A four-part evening series through the Hitchhiker's Guide to Agentic AI, Jun 30 – Jul 14."]):
        _rounded(s, GX + Inches(4.45), Inches(3.05 + i * 0.85), Inches(3.33), Inches(0.72), PAPER_SOFT, INK, 1.0, False, 0.12)
        tb(s, GX + Inches(4.6), Inches(3.1 + i * 0.85), Inches(3.05), Inches(0.62), t, 8.5, False, INK_SOFT, BODY)
    # col 3 saturday (tone-green)
    _rounded(s, GX + Inches(8.3), Inches(2.5), Inches(3.93), Inches(3.4), T_GREEN, INK, 1.5, True, 0.08)
    tb(s, GX + Inches(8.6), Inches(2.65), Inches(3.3), Inches(0.3), "▦  WHAT A SATURDAY LOOKS LIKE", 8, True, INK_SOFT, BODY)
    for i, (a, b2) in enumerate([("9:00 — Doors & check-in", "RSVP plus the event Google form at the LSEG gate, then laptops open."),
                                  ("Morning — Build, not slides", "Live coding you follow along: tensors, attention, agent loops."),
                                  ("Close — Verify & network", "Demos, Q&A, and the hallway track where teams form.")]):
        tb(s, GX + Inches(8.9), Inches(3.1 + i * 0.85), Inches(2.9), Inches(0.3), f"{i+1}  {a}", 9, True, INK, BODY)
        tb(s, GX + Inches(8.9), Inches(3.35 + i * 0.85), Inches(2.9), Inches(0.5), b2, 8.5, False, INK_SOFT, BODY)
    footer(s)

    # -- 4. Approach (ink dark, exact) --
    s = prs.slides.add_slide(blank); bg(s, INK); header(s, dark=True)
    tb(s, GX, Inches(1.0), GW, Inches(0.3), "OUR APPROACH", 9, True, T_YELLOW, BODY)
    tb(s, GX, Inches(1.3), Inches(8.0), Inches(1.2), "First you research. Then you build. Then you transform.", 30, True, PAPER_SOFT, HEAD)
    tb(s, GX, Inches(2.35), Inches(8.0), Inches(0.5), "Yatra means journey. Ours takes you from curious onlooker to confident builder — through research, build, and transform, repeated every single meetup.", 11, False, PAPER_SOFT, BODY)
    steps = [
        ("01 · Research", "Research", "We read the papers, trace the math, and ask the naive questions out loud — so nobody has to pretend they already know.", T_VIOLET),
        ("02 · Build", "Build", "Laptops open, code on screen. Every session ships something real — a training loop, an agent harness, a working demo.", T_GREEN),
        ("03 · Transform", "Transform", "Skills become careers, side projects become products, and strangers become collaborators. That is the AIYatra journey.", T_CORAL),
    ]
    for i, (kick, title, body, tone) in enumerate(steps):
        x = GX + Inches(i * 4.15)
        _rounded(s, x, Inches(3.15), Inches(3.93), Inches(2.4), PAPER_SOFT, PAPER_SOFT, 1.0, False, 0.08)
        pill(s, x + Inches(0.3), Inches(3.35), Inches(1.9), Inches(0.34), kick, tone)
        tb(s, x + Inches(0.3), Inches(3.8), Inches(3.3), Inches(0.5), title, 26, True, INK, HEAD)
        tb(s, x + Inches(0.3), Inches(4.35), Inches(3.3), Inches(1.0), body, 10, False, INK_SOFT, BODY)
    tb(s, GX, Inches(5.85), GW, Inches(0.3), "14 Meetups hosted   •   2hr Per session   •   100% Free, always", 9, True, T_YELLOW, BODY, PP_ALIGN.CENTER)
    footer(s)

    # -- 5. Testimonials (paper) --
    s = prs.slides.add_slide(blank); bg(s, PAPER); header(s)
    kicker_line(s, Inches(1.0), "TESTIMONIALS")
    h_hand(s, Inches(1.28), "Rated 4.6★ by the people who show up.", 32)
    tb(s, GX + Inches(9.2), Inches(1.15), Inches(3.0), Inches(0.4), "★  4.6 · 90 ratings on Meetup", 9, True, INK, BODY)
    quotes = [
        ("I walked in knowing nothing about transformers and walked out having built attention from scratch. No gatekeeping, no jargon walls — just patient, brilliant teaching.", "Priya S.", "Data Analyst → ML Engineer", T_BLUE),
        ("The agentic AI workshop was the best Saturday I have spent in years. We built a coding agent on our own machines, no API keys, and it actually worked.", "Rahul K.", "Backend Engineer", T_GREEN),
        ("As a student, most AI events felt out of reach. AIYatra is free, welcoming, and genuinely deep — the linear algebra session finally made the math click for me.", "Sai Rishita M.", "CS Undergraduate", T_YELLOW),
        ("You come for the sessions and stay for the people. I found my co-founder, my study group, and my confidence here. This community changes trajectories.", "Kiran K.", "Founder, AI Startup", T_CORAL),
    ]
    for i, (q, n, r, tone) in enumerate(quotes):
        x = GX + Inches(i * 3.1)
        _rounded(s, x, Inches(2.3), Inches(2.9), Inches(3.4), tone, INK, 1.5, True, 0.08)
        tb(s, x + Inches(0.25), Inches(2.55), Inches(2.4), Inches(2.0), f'“{q}”', 9.5, False, INK, BODY)
        tb(s, x + Inches(0.25), Inches(4.85), Inches(2.4), Inches(0.3), n, 9.5, True, INK, BODY)
        tb(s, x + Inches(0.25), Inches(5.1), Inches(2.4), Inches(0.3), r.upper(), 7.5, True, INK_SOFT, BODY)
    footer(s)

    # -- 6. Ambassador band (events bg, exact) --
    s = prs.slides.add_slide(blank); bg(s, EVENTS_BG); header(s)
    _rounded(s, GX, Inches(1.4), GW, Inches(3.6), PAPER_SOFT, INK, 1.75, True, 0.08)
    try:
        s.shapes.add_picture(CREST, GX + Inches(0.4), Inches(1.8), Inches(2.2), Inches(2.2))
    except: pass
    tb(s, GX + Inches(3.0), Inches(1.7), Inches(5.5), Inches(0.3), "NEW CHAPTER · STUDENT AMBASSADOR PROGRAM", 8.5, True, INK_SOFT, BODY)
    tb(s, GX + Inches(3.0), Inches(2.0), Inches(5.5), Inches(0.8), "Carry the yatra to your campus.", 32, True, INK, HEAD)
    tb(s, GX + Inches(3.0), Inches(2.9), Inches(5.5), Inches(0.7), "Rally your batch to Saturday meetups, host campus mini-sessions, and earn recognition, mentorship and speaker slots — free, always.", 11, False, INK_SOFT, BODY)
    btn(s, GX + Inches(3.0), Inches(3.85), Inches(2.7), Inches(0.52), "Explore the program  →")
    # organizers strip (exact tiles, compact)
    tb(s, GX, Inches(5.35), GW, Inches(0.35), "The people behind AIYatra — Khaja Moinuddin Mohammed (Super Organizer) · Azeez Syed (Co-organizer & Host) · Jagadeeswara Reddy (Host & Educator)", 9, True, INK, BODY, PP_ALIGN.CENTER)
    footer(s)

    # -- 7. About (about green tint, exact) --
    s = prs.slides.add_slide(blank); bg(s, ABOUT_BG); header(s)
    kicker_line(s, Inches(1.0), "ABOUT AIYATRA")
    tb(s, GX, Inches(1.28), Inches(6.4), Inches(1.1), "An open door into AI, in the heart of Hyderabad.", 32, True, INK, HEAD)
    tb(s, GX, Inches(2.3), Inches(6.4), Inches(1.3),
       "AIYatra is for anyone exploring AI applications, diving into Artificial Intelligence: A Modern Approach, building recommender systems, or mastering machine learning with Python. Beginner or expert — you get knowledge sharing, networking, and collaborative projects.\nLed by super organizer Khaja Moinuddin Mohammed and a crew of passionate volunteers, we meet at LSEG, International Tech Park, Madhapur, Hyderabad.",
       10.5, False, INK_SOFT, BODY)
    _rounded(s, GX + Inches(7.0), Inches(1.3), Inches(5.23), Inches(4.3), PAPER_SOFT, INK, 1.5, True, 0.08)
    tb(s, GX + Inches(7.4), Inches(1.55), Inches(4.4), Inches(0.4), "The essentials", 24, True, INK, HEAD)
    for i, (k, v) in enumerate([("Where", "LSEG, International Tech Park, Madhapur, Hyderabad"), ("When", "Saturdays · mornings, IST"), ("Cost", "Free, always — learning should be"), ("Bring", "A laptop, Python 3.10+, and curiosity")]):
        _rounded(s, GX + Inches(7.4), Inches(2.15 + i * 0.82), Inches(4.43), Inches(0.68), PAPER, INK, 1.0, False, 0.12)
        tb(s, GX + Inches(7.6), Inches(2.18 + i * 0.82), Inches(4.1), Inches(0.22), k.upper(), 7.5, True, INK_SOFT, BODY)
        tb(s, GX + Inches(7.6), Inches(2.38 + i * 0.82), Inches(4.1), Inches(0.3), v, 9.5, True, INK, BODY)
    footer(s)

    # -- 8. Final CTA (ink dark, exact) --
    s = prs.slides.add_slide(blank); bg(s, INK); header(s, dark=True)
    pill(s, Inches(5.15), Inches(1.3), Inches(3.0), Inches(0.38), "🎟  Free forever · No prerequisites")
    tb(s, GX, Inches(1.85), GW, Inches(1.3), "Your AI journey starts with a single RSVP.", 44, True, PAPER_SOFT, HEAD, PP_ALIGN.CENTER)
    tb(s, Inches(2.8), Inches(3.0), Inches(7.7), Inches(0.5), "Join 2,957 learners democratizing AI — one meetup, one project, one breakthrough at a time.", 12, False, PAPER_SOFT, BODY, PP_ALIGN.CENTER)
    btn(s, Inches(3.6), Inches(3.8), Inches(2.9), Inches(0.55), "Join AIYatra on Meetup  →", T_YELLOW, INK)
    _rounded(s, Inches(6.8), Inches(3.8), Inches(2.9), Inches(0.55), INK, PAPER_SOFT, 1.5, False, 0.18)
    tb(s, Inches(6.8), Inches(3.8), Inches(2.9), Inches(0.55), "RSVP: Agentic AI, Sep 5", 10.5, True, PAPER_SOFT, BODY, PP_ALIGN.CENTER).text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    footer(s)

    # -- 9. Contact / footer replica (ink) --
    s = prs.slides.add_slide(blank); bg(s, INK)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SW, Pt(6))
    top.fill.solid(); top.fill.fore_color.rgb = T_YELLOW; top.line.fill.background()
    header(s, dark=True)
    tb(s, GX, Inches(1.1), GW, Inches(0.6), "AI Yatra — Research. Build. Transform.", 30, True, PAPER_SOFT, HEAD)
    tb(s, GX, Inches(1.7), Inches(7.0), Inches(0.4), "Hyderabad's open AI community. We research, build, and transform — one hands-on Saturday at a time. Free, forever, and open to everyone.", 10.5, False, PAPER_SOFT, BODY)
    cols = [
        ("EXPLORE", "Ambassadors — new\nBlog\nEvents\nCommunity\nTestimonials"),
        ("SHOW UP", "LSEG, International Tech Park, Madhapur\nSaturdays · mornings, IST\nFree, always"),
        ("FOLLOW", f"{EMAIL}\nMeetup inbox\nGitHub: {GITHUB_SHORT}\nLinkedIn: {LINKEDIN_SHORT}"),
        ("VISIT", f"{WEBSITE}\n{WEBSITE}/ambassadors\n{WEBSITE}/blog\n{MEETUP_SHORT}"),
    ]
    for i, (h, body) in enumerate(cols):
        x = GX + Inches(i * 3.1)
        tb(s, x, Inches(2.6), Inches(2.9), Inches(0.3), h, 9, True, T_YELLOW, BODY)
        tb(s, x, Inches(2.95), Inches(2.9), Inches(1.8), body, 9.5, False, PAPER_SOFT, BODY)
    _rounded(s, GX, Inches(5.0), GW, Inches(1.15), PAPER_SOFT, PAPER_SOFT, 1.0, False, 0.1)
    tb(s, GX + Inches(0.4), Inches(5.15), Inches(11.4), Inches(0.85),
       f"Website: {WEBSITE}   •   LinkedIn: {LINKEDIN}   •   Email: {EMAIL}   •   Meetup: {MEETUP}   •   Made with ♥ in Hyderabad   •   © 2026 AIYatra",
       9.5, True, INK, BODY, PP_ALIGN.CENTER).text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    footer(s)

    p1 = OUT_DIR + "/AIYatra-Journey-From-Start-Till-Now.pptx"
    prs.save(p1); print("Saved:", p1)
    return p1

# ================= DECK 2 =================
def deck2():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    blank = prs.slide_layouts[6]

    # 1 Hero (exact AmbassadorsPage hero)
    s = prs.slides.add_slide(blank); bg(s, PAPER); header(s)
    pill(s, GX, Inches(1.0), Inches(4.2), Inches(0.38), "🎓  New chapter · Student Ambassador Program", T_BLUE)
    tb(s, GX, Inches(1.5), Inches(7.0), Inches(1.5), "Carry the AI Yatra to your campus.", 44, True, INK, HEAD)
    tb(s, GX, Inches(2.7), Inches(6.6), Inches(0.8),
       "AIYatra Student Ambassadors are the face of AIYatra in their colleges — rallying classmates to Saturday meetups, hosting campus mini-sessions, and turning dense AI topics into hands-on learning their batch can actually use.",
       11.5, False, INK_SOFT, BODY)
    btn(s, GX, Inches(3.75), Inches(2.4), Inches(0.52), "Apply by email  →")
    _rounded(s, GX + Inches(2.6), Inches(3.75), Inches(2.7), Inches(0.52), PAPER_SOFT, INK, 1.5, True, 0.18)
    tb(s, GX + Inches(2.6), Inches(3.75), Inches(2.7), Inches(0.52), "Attend a meetup first", 10.5, True, INK, BODY, PP_ALIGN.CENTER).text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _rounded(s, Inches(7.9), Inches(1.0), Inches(4.85), Inches(4.6), PAPER_SOFT, INK, 1.75, True, 0.08)
    try:
        s.shapes.add_picture(CREST, Inches(8.1), Inches(1.2), Inches(4.45), Inches(3.4))
    except: pass
    tb(s, Inches(8.1), Inches(4.75), Inches(4.45), Inches(0.6), f"Apply: {EMAIL}  •  {WEBSITE_SHORT}/ambassadors", 8.5, True, T_BLUE_D, BODY, PP_ALIGN.CENTER)
    footer(s)

    # 2 Mandate (paper-soft, exact duties)
    s = prs.slides.add_slide(blank); bg(s, PAPER_SOFT); header(s)
    kicker_line(s, Inches(1.0), "✦  THE MANDATE")
    h_hand(s, Inches(1.28), "Six duties. One community.", 32)
    tb(s, GX, Inches(1.9), GW, Inches(0.35), "A term runs one semester. Hit these six notes and you renew with honors — and a letter that actually says something.", 10.5, False, INK_SOFT, BODY)
    duties = [
        ("Carry the AI Yatra to campus", "Spread the word in your college — posters, classroom shout-outs, WhatsApp groups and tech clubs. You are the bridge between AIYatra and your campus.", T_YELLOW),
        ("Bring your batch along", "Rally classmates to Saturday meetups, help first-timers settle in, and keep the newcomer energy warm from RSVP to after-party notes.", T_GREEN),
        ("Host campus mini-sessions", "Run paper-readings, build nights and revision circles at your college with our curriculum kits — same research·build·transform rhythm, your turf.", T_BLUE),
        ("Document the journey", "Capture photos, clips and field notes from meetups and campus sessions so the community story keeps growing online.", T_CORAL),
        ("Learn in the fast lane", "Get first access to session material, prep calls with organizers, and guided tracks — transformers, PyTorch, agentic AI — before anyone else.", T_VIOLET),
        ("Feed the roadmap back", "Tell us what students actually want: topics, timings, formats. Ambassador feedback directly shapes the upcoming meetup arc.", T_YELLOW),
    ]
    for i, (t, b2, tone) in enumerate(duties):
        x = GX + Inches((i % 3) * 4.15); y = Inches(2.45 + (i // 3) * 2.0)
        _rounded(s, x, y, Inches(3.93), Inches(1.8), PAPER, INK, 1.5, True, 0.08)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.25), Pt(14), Pt(14))
        dot.fill.solid(); dot.fill.fore_color.rgb = tone; dot.line.color.rgb = INK; dot.line.width = Pt(1.25)
        tb(s, x + Inches(0.6), y + Inches(0.2), Inches(3.0), Inches(0.45), t, 13, True, INK, HEAD)
        tb(s, x + Inches(0.3), y + Inches(0.75), Inches(3.33), Inches(0.9), b2, 8.5, False, INK_SOFT, BODY)
    footer(s)

    # 3 Perks (ink dark, exact)
    s = prs.slides.add_slide(blank); bg(s, INK); header(s, dark=True)
    tb(s, GX, Inches(1.0), GW, Inches(0.3), "WHAT YOU GET", 9, True, T_YELLOW, BODY)
    tb(s, GX, Inches(1.3), GW, Inches(0.9), "Ambassadorship pays in proof of work.", 32, True, PAPER_SOFT, HEAD)
    tb(s, GX, Inches(2.0), Inches(7.5), Inches(0.4), "No stipends, no swag-bribes — everything here compounds into skills, relationships, and receipts future employers can verify.", 10.5, False, PAPER_SOFT, BODY)
    perks = [
        ("Official recognition", "A signed Ambassador certificate plus a public profile on this page once you complete your first term."),
        ("Letters that count", "Recommendation letters for internships and higher studies, earned through real shipped work — not attendance."),
        ("Direct mentor access", "Monthly office hours with the organizers. Bring your projects, career questions, and wildest paper ideas."),
        ("Speaker lane", "Ambassadors in good standing get priority slots to demo, lightning-talk, and eventually lead sessions."),
    ]
    for i, (t, b2) in enumerate(perks):
        x = GX + Inches((i % 2) * 6.2); y = Inches(2.7 + (i // 2) * 1.65)
        _rounded(s, x, y, Inches(6.03), Inches(1.4), PAPER_SOFT, PAPER_SOFT, 1.0, False, 0.08)
        _rounded(s, x + Inches(0.3), y + Inches(0.3), Inches(0.55), Inches(0.55), T_YELLOW, INK, 1.25, False, 0.2)
        tb(s, x + Inches(0.35), y + Inches(0.35), Inches(0.5), Inches(0.45), "✓", 14, True, INK, BODY, PP_ALIGN.CENTER)
        tb(s, x + Inches(1.05), y + Inches(0.25), Inches(4.6), Inches(0.4), t, 19, True, INK, HEAD)
        tb(s, x + Inches(1.05), y + Inches(0.7), Inches(4.6), Inches(0.55), b2, 9.5, False, INK_SOFT, BODY)
    footer(s)

    # 4 Who fits + How to apply (green tint, exact two-col)
    s = prs.slides.add_slide(blank); bg(s, ABOUT_BG); header(s)
    _rounded(s, GX, Inches(1.05), Inches(6.0), Inches(4.95), PAPER_SOFT, INK, 1.5, True, 0.08)
    tb(s, GX + Inches(0.4), Inches(1.25), Inches(5.2), Inches(0.3), "WHO FITS", 9, True, T_BLUE_D, BODY)
    tb(s, GX + Inches(0.4), Inches(1.5), Inches(5.2), Inches(0.7), "Built for students who show up.", 26, True, INK, HEAD)
    fits = [
        "Enrolled in any college or university — any year, any branch. Curiosity beats CGPA here.",
        "Attended at least one AIYatra meetup (or will attend the next one) — you should know the room you are inviting people into.",
        "Can give ~3 hours a week: one meetup plus campus outreach and a short weekly check-in.",
        "Comfortable talking to people — in classrooms, clubs, and group chats — and following through on what you promise.",
    ]
    box = tb(s, GX + Inches(0.4), Inches(2.35), Inches(5.2), Inches(2.2), "", 9, False, INK, BODY)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "①  " + fits[0]; p.font.size = Pt(9); p.font.name = BODY; p.font.color.rgb = INK; p.space_after = Pt(6)
    for j, f in enumerate(fits[1:], 2):
        num = ["②", "③", "④"][j-2]
        para(tf, f"{num}  " + f, 9, False, INK, BODY, PP_ALIGN.LEFT, Pt(6))
    _rounded(s, GX + Inches(0.4), Inches(4.85), Inches(5.2), Inches(0.9), T_YELLOW, INK, 1.25, False, 0.12)
    tb(s, GX + Inches(0.6), Inches(4.9), Inches(4.8), Inches(0.3), "♥  Missing one? Apply anyway.", 11, True, INK, HEAD)
    tb(s, GX + Inches(0.6), Inches(5.2), Inches(4.8), Inches(0.4), "Curiosity and follow-through beat credentials — most ambassadors started as first-timers in the back row.", 8.5, False, INK_SOFT, BODY)
    # right col
    _rounded(s, GX + Inches(6.23), Inches(1.05), Inches(6.0), Inches(4.95), PAPER, INK, 1.5, True, 0.08)
    tb(s, GX + Inches(6.63), Inches(1.25), Inches(5.2), Inches(0.3), "HOW TO APPLY", 9, True, T_BLUE_D, BODY)
    tb(s, GX + Inches(6.63), Inches(1.5), Inches(5.2), Inches(0.7), "Three steps to the crest.", 26, True, INK, HEAD)
    for i, (n, t, b2) in enumerate([("01", "Say hello", "Email us with your name, college, year, and one paragraph on why you want to carry the AI Yatra to your campus."),
                                     ("02", "Meet the crew", "A short intro call with an organizer — we understand your campus scene and agree on your first 30-day plan."),
                                     ("03", "Run your first drive", "Bring five classmates to a Saturday meetup and host one campus huddle. Do that and the crest is yours.")]):
        _rounded(s, GX + Inches(6.63), Inches(2.35 + i * 1.0), Inches(5.2), Inches(0.88), PAPER_SOFT, INK, 1.0, False, 0.12)
        tb(s, GX + Inches(6.8), Inches(2.4 + i * 1.0), Inches(4.9), Inches(0.25), f"{n}  ·  {t}", 10, True, T_BLUE_D, BODY)
        tb(s, GX + Inches(6.8), Inches(2.62 + i * 1.0), Inches(4.9), Inches(0.5), b2, 8.5, False, INK_SOFT, BODY)
    tb(s, GX + Inches(6.63), Inches(5.45), Inches(5.2), Inches(0.35), f"Apply at {EMAIL}  •  Open to every college · Free, always", 8, True, INK, BODY, PP_ALIGN.CENTER)
    footer(s)

    # 5 Closing CTA (ink, exact)
    s = prs.slides.add_slide(blank); bg(s, INK); header(s, dark=True)
    tb(s, GX, Inches(1.3), GW, Inches(0.4), "ONE CAMPUS AT A TIME", 10, True, T_YELLOW, BODY, PP_ALIGN.CENTER)
    tb(s, GX, Inches(1.7), GW, Inches(1.3), "Your campus is the next AI Yatra.", 46, True, PAPER_SOFT, HEAD, PP_ALIGN.CENTER)
    tb(s, Inches(3.0), Inches(2.9), Inches(7.3), Inches(0.4), "Write to us today — or come experience a Saturday first, then decide.", 12, False, PAPER_SOFT, BODY, PP_ALIGN.CENTER)
    btn(s, Inches(3.6), Inches(3.6), Inches(2.9), Inches(0.55), "Become an ambassador  →", T_YELLOW, INK)
    _rounded(s, Inches(6.8), Inches(3.6), Inches(2.9), Inches(0.55), INK, PAPER_SOFT, 1.5, False, 0.18)
    tb(s, Inches(6.8), Inches(3.6), Inches(2.9), Inches(0.55), "Join on Meetup ↗", 10.5, True, PAPER_SOFT, BODY, PP_ALIGN.CENTER).text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb(s, GX, Inches(4.6), GW, Inches(0.4), f"{EMAIL}   •   {WEBSITE_SHORT}/ambassadors   •   {MEETUP_SHORT}", 9.5, True, T_YELLOW, BODY, PP_ALIGN.CENTER)
    footer(s)

    # 6 Contact (ink footer replica)
    s = prs.slides.add_slide(blank); bg(s, INK)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SW, Pt(6))
    top.fill.solid(); top.fill.fore_color.rgb = T_YELLOW; top.line.fill.background()
    header(s, dark=True)
    tb(s, GX, Inches(1.1), GW, Inches(0.6), "Student Ambassadors — start here.", 30, True, PAPER_SOFT, HEAD)
    tb(s, GX, Inches(1.7), Inches(8.0), Inches(0.4), "One email is all it takes. Tell us your name, college, year, and why you want to carry the AI Yatra to your campus.", 10.5, False, PAPER_SOFT, BODY)
    items = [("EMAIL", EMAIL, "Application + questions"), ("WEB", f"{WEBSITE}/ambassadors", "Program details"), ("LINKEDIN", LINKEDIN_SHORT, "Follow + milestones"), ("MEETUP", MEETUP_SHORT, "Attend first: Sep 5")]
    for i, (k, v, d) in enumerate(items):
        x = GX + Inches(i * 3.1)
        _rounded(s, x, Inches(2.5), Inches(2.9), Inches(1.9), PAPER_SOFT, PAPER_SOFT, 1.0, False, 0.08)
        tb(s, x + Inches(0.2), Inches(2.7), Inches(2.5), Inches(0.3), k, 9, True, T_BLUE_D, BODY, PP_ALIGN.CENTER)
        tb(s, x + Inches(0.2), Inches(3.0), Inches(2.5), Inches(0.6), v, 9.5, True, INK, BODY, PP_ALIGN.CENTER)
        tb(s, x + Inches(0.2), Inches(3.6), Inches(2.5), Inches(0.4), d, 9, False, INK_SOFT, BODY, PP_ALIGN.CENTER)
    _rounded(s, GX, Inches(4.8), GW, Inches(1.15), PAPER_SOFT, PAPER_SOFT, 1.0, False, 0.1)
    tb(s, GX + Inches(0.4), Inches(4.95), Inches(11.4), Inches(0.85),
       f"Website: {WEBSITE}   •   LinkedIn: {LINKEDIN}   •   Email: {EMAIL}   •   Meetup: {MEETUP}   •   © 2026 AIYatra · Research. Build. Transform.",
       9.5, True, INK, BODY, PP_ALIGN.CENTER).text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    footer(s)

    p2 = OUT_DIR + "/AIYatra-Student-Ambassador-Program.pptx"
    prs.save(p2); print("Saved:", p2)
    return p2

if __name__ == "__main__":
    deck1(); deck2(); print("Done v2.")
